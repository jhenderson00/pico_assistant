import openai
import csv
import PyPDF2
import numpy as np
from openai.embeddings_utils import cosine_similarity
from scipy.spatial import distance_matrix
import docx
from striprtf.striprtf import rtf_to_text
from odf import text, teletype
from odf.opendocument import load
from pptx import Presentation
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import warnings

openai.api_key = open("openai_key.txt", "r").read().strip("\n")  # get api key from text file

#Takes a path to a PDF and returns the text contents
def pdf_to_text(file_path):
    pdf_reader = PyPDF2.PdfReader(file_path)
    text = ""
    for page_num in range(len(pdf_reader.pages)):
        page_obj = pdf_reader.pages[page_num]
        pages = page_obj.extract_text()
        text += pages
    return text

def read_word_file(file_path):
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def read_txt_file(file_path):
    with open(file_path, 'r') as file:
        text = file.read()
    return text

def read_rtf_file(file_path):
    with open(file_path) as infile:
        content = infile.read()
        text = rtf_to_text(content)
    return text

def read_odt_file(file_path):
    full_text = []
    odt_doc = load(file_path)
    paragraphs = odt_doc.getElementsByType(text.P)
    for i in paragraphs:
        paragraph = teletype.extractText(i)
        full_text.append(paragraph)
    return '\n'.join(full_text)

def read_ppt_file(file_path):
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Split text into words and join with a single space to remove extra whitespace
                    text = ' '.join(paragraph.text.split())
                    # Only append non-empty text
                    if text.strip():
                        full_text.append(text)
    # Join paragraphs with a single newline
    return '\n'.join(full_text)

def read_epub_file(file_path):
    # Filter out the ebooklib warning
    with warnings.catch_warnings():
        warnings.filterwarnings("ignore", category=UserWarning)
        book = epub.read_epub(file_path)
    full_text = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), 'html.parser')
        text = soup.get_text()
        # Split text into words and join with a single space to remove extra whitespace
        cleaned_text = ' '.join(text.split())
        if cleaned_text.strip():
            full_text.append(cleaned_text)
    return '\n'.join(full_text)

#Split the input text into smaller chunks of a specified size.
def split_text(text, chunk_size):
    text_chunks = []
    text_length = len(text)
    start = 0
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        text_chunks.append(chunk)
        start = end
    return text_chunks

def create_embeddings(text_chunks, model="text-embedding-ada-002"):
    embeddings = []
    try:
        prepared_chunks = [chunk.replace("\n", " ") for chunk in text_chunks]
        response = openai.Embedding.create(input=prepared_chunks, model=model)
        if response and "data" in response:
            for data in response["data"]:
                embeddings.append(data["embedding"])
        return embeddings
    except Exception as e:
        print(f"Error creating embeddings: {e}")
        return None

def write_embeddings_to_csv(embeddings, csv_path):
    with open(csv_path, "w", newline="") as csvfile:
        csv_writer = csv.writer(csvfile)
        for embedding in embeddings:
            csv_writer.writerow(embedding)

def read_embeddings_from_csv(csv_path):
    embeddings = []
    with open(csv_path, "r", newline="") as csvfile:
        csv_reader = csv.reader(csvfile)
        for row in csv_reader:
            embedding = [float(value) for value in row]
            embeddings.append(embedding)
    return embeddings

def calculate_centroid(embeddings):
    centroid = np.mean(embeddings, axis=0)
    return centroid

def closest_embeddings_to_centroid(embeddings, centroid, n=3):
    distances = [distance_matrix([embedding], [centroid])[0][0] for embedding in embeddings]
    closest_indices = np.argpartition(distances, range(n))[:n]
    return closest_indices.tolist()

def search_embeddings(query, embeddings, n=3):
    """
    Search for the most similar embeddings to the given query using cosine similarity.

    Args:
        query (str): The input query.
        embeddings (list): A list of embedding vectors.
        n (int): The number of top results to return.

    Returns:
        list: A list of indices of the top N most similar embeddings.
    """
    query_embedding = create_embeddings([query])[0]
    similarities = [cosine_similarity(embedding, query_embedding) for embedding in embeddings]

    # Get the indices of the top N most similar embeddings
    top_indices = np.argsort(similarities)[-n:][::-1]

    return top_indices.tolist()

def retrieve_answer(indices, text_chunks, n=1):
    """
    Retrieve the most relevant text from the text chunks using the provided indices.

    Args:
        indices (list): A list of indices of the most similar embeddings.
        text_chunks (list): A list of text chunks.
        n (int): The number of top answers to return.

    Returns:
        list: A list of the top N most relevant answers.
    """
    if n > len(indices):
        n = len(indices)
    answers = [text_chunks[index] for index in indices[:n]]
    return answers

def summarize_text(embeddings, text_chunks, n=3):
    centroid = calculate_centroid(embeddings)
    closest_indices = closest_embeddings_to_centroid(embeddings, centroid, n)
    summary = retrieve_answer(closest_indices, text_chunks, n)
    return summary

def process_pdfs_and_create_csv(pdf_paths, csv_path, chunk_size=1000):
    all_chunks = []
    all_embeddings = []
    for pdf_path in pdf_paths:
        text = pdf_to_text(pdf_path)
        chunks = split_text(text, chunk_size)
        embeddings = create_embeddings(chunks)
        all_chunks.extend(chunks)
        all_embeddings.extend(embeddings)
    write_embeddings_to_csv(all_embeddings, csv_path)
    return csv_path, all_chunks

#pdf_paths = ["./pdfs/living_in_the_light.pdf"]
#csv_path, text_chunks = process_pdfs_and_create_csv(pdf_paths, "./pdfs/living_in_the_light.csv", chunk_size=500)
#saved_embeds = read_embeddings_from_csv(csv_path)
#summary_chunks = summarize_text(saved_embeds, text_chunks, n=3)
#print(summary_chunks)

print(read_epub_file("./pdfs/20000-Leagues-Under-the-Sea.epub"))




